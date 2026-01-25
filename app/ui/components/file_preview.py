import os
import json
from pathlib import Path
from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QLabel, QScrollArea, QTextEdit, 
    QFrame, QHBoxLayout, QPushButton, QSizePolicy
)
from PySide6.QtCore import Qt, QSize, Signal
from PySide6.QtGui import QPixmap, QFont, QIcon

from core.config import config
from ai.gemini_client import analyze_file_with_gemini, analyze_text_with_gemini

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

# Office document support
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Legacy Office document support
try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    import olefile
    OLEFILE_AVAILABLE = True
except ImportError:
    OLEFILE_AVAILABLE = False

try:
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False


class FilePreview(QWidget):
    """Component to preview different types of documents"""
    
    # Signal emitted when preview content changes
    content_changed = Signal(str)  # file_path
    
    # Signals emitted when user wants to send text to form fields
    send_to_date_requested = Signal(str)  # selected_text
    send_to_organization_requested = Signal(str)  # selected_text
    send_to_subject_requested = Signal(str)  # selected_text
    send_to_receiver_requested = Signal(str)  # selected_text
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.current_file_path = None
        self.current_pdf_doc = None  # Store PDF document for navigation
        self.current_page_num = 0  # Current page number (0-based)
        self.total_pages = 0  # Total number of pages
        self._extracted_text_content = None  # Store extracted text (OCR or AI)
        self._is_showing_extracted_text = False  # Flag to prevent re-rendering
        self._ai_result_cache = {}  # Cache AI results by file path
        self._current_ai_header = None  # Current AI header text
        self._setup_ui()
    
    def _setup_ui(self):
        """Setup the preview UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(10, 10, 10, 10)
        
        # Set size policy to expand in both directions
        self.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        # Header with file info
        self._setup_header(layout)
        
        # Preview content area
        self._setup_content_area(layout)
        
        # Initially show placeholder
        self._show_placeholder()
    
    def _setup_header(self, parent_layout):
        """Setup the header section with file info"""
        header_frame = QFrame()
        header_frame.setFrameStyle(QFrame.Box)
        header_frame.setMaximumHeight(120)  # Increased to accommodate navigation buttons
        
        header_layout = QVBoxLayout(header_frame)
        
        # File name label
        self._file_name_label = QLabel(self.tr("No document selected"))
        font = QFont()
        font.setBold(True)
        font.setPointSize(12)
        self._file_name_label.setFont(font)
        self._file_name_label.setWordWrap(True)
        header_layout.addWidget(self._file_name_label)
        
        # File info label
        self._file_info_label = QLabel("")
        self._file_info_label.setStyleSheet("color: #666;")
        header_layout.addWidget(self._file_info_label)
        
        # PDF navigation controls
        self._setup_pdf_navigation(header_layout)
        
        parent_layout.addWidget(header_frame)
    
    def _setup_pdf_navigation(self, parent_layout):
        """Setup PDF page navigation controls"""
        self._nav_widget = QWidget()
        nav_layout = QHBoxLayout(self._nav_widget)
        nav_layout.setContentsMargins(0, 5, 0, 0)
        
        # Previous page button
        self._prev_page_btn = QPushButton("◀ " + self.tr("Previous"))
        self._prev_page_btn.setMaximumHeight(30)
        self._prev_page_btn.clicked.connect(self._go_to_previous_page)
        self._prev_page_btn.setEnabled(False)
        nav_layout.addWidget(self._prev_page_btn)
        
        # Page info label
        self._page_info_label = QLabel("")
        self._page_info_label.setAlignment(Qt.AlignCenter)
        self._page_info_label.setStyleSheet("font-weight: bold; color: #333;")
        nav_layout.addWidget(self._page_info_label)
        
        # Next page button
        self._next_page_btn = QPushButton(self.tr("Next") + " ▶")
        self._next_page_btn.setMaximumHeight(30)
        self._next_page_btn.clicked.connect(self._go_to_next_page)
        self._next_page_btn.setEnabled(False)
        nav_layout.addWidget(self._next_page_btn)
        
        # OCR button
        self._ocr_btn = QPushButton("🔤 " + self.tr("Extract Text"))
        self._ocr_btn.setMaximumHeight(30)
        self._ocr_btn.setStyleSheet("""
            QPushButton {
                background-color: #007acc;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 5px 10px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #005a9e;
            }
            QPushButton:pressed {
                background-color: #004578;
            }
        """)
        self._ocr_btn.clicked.connect(self._extract_pdf_text)
        self._ocr_btn.setEnabled(False)
        nav_layout.addWidget(self._ocr_btn)

        # AI Analyze button
        self._ai_btn = QPushButton("🤖 " + self.tr("Analyze with AI"))
        self._ai_btn.setMaximumHeight(30)
        self._ai_btn.setStyleSheet("""
            QPushButton {
                background-color: #6f42c1;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 5px 10px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #59359c;
            }
            QPushButton:pressed {
                background-color: #4b2d82;
            }
            QPushButton:disabled {
                background-color: #b9a5e6;
                color: #f1ecff;
            }
        """)
        self._ai_btn.clicked.connect(self._analyze_with_ai)
        self._ai_btn.setEnabled(False)
        self._ai_btn.setVisible(False)
        nav_layout.addWidget(self._ai_btn)
        
        # Initially hide navigation controls
        self._nav_widget.setVisible(False)
        parent_layout.addWidget(self._nav_widget)
    
    def _setup_content_area(self, parent_layout):
        """Setup the scrollable content area"""
        # Create scroll area
        self._scroll_area = QScrollArea()
        self._scroll_area.setWidgetResizable(True)
        self._scroll_area.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self._scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self._scroll_area.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        # Content widget inside scroll area
        self._content_widget = QWidget()
        self._content_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self._content_layout = QVBoxLayout(self._content_widget)
        
        # Preview label/widget
        self._preview_widget = QLabel()
        self._preview_widget.setAlignment(Qt.AlignCenter)
        self._preview_widget.setMinimumSize(300, 400)
        self._preview_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self._preview_widget.setScaledContents(False)  # Important for images
        self._content_layout.addWidget(self._preview_widget)
        
        self._scroll_area.setWidget(self._content_widget)
        parent_layout.addWidget(self._scroll_area)
    
    def preview_file(self, file_path):
        """Preview the given file"""
        if not file_path or not os.path.exists(file_path):
            self._show_error(self.tr("File not found"))
            return
        
        self.current_file_path = file_path
        file_info = Path(file_path)
        
        # Update header
        self._file_name_label.setText(file_info.name)
        
        # Get file size and extension
        file_size = self._format_file_size(file_info.stat().st_size)
        file_ext = file_info.suffix.upper()
        self._file_info_label.setText(f"{file_ext} • {file_size}")
        
        # Preview based on file type
        self._preview_by_type(file_path, file_info.suffix.lower())
        
        # Emit signal
        self.content_changed.emit(file_path)
    
    def _preview_by_type(self, file_path, extension):
        """Preview file based on its type"""
        # Reset any previous state (OCR buttons, etc.)
        self._reset_preview_state()
        
        if extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif']:
            self._preview_image(file_path)
            self._show_ai_for_simple_preview()
        elif extension in ['.txt', '.rtf']:
            self._preview_text(file_path)
            self._show_ai_for_simple_preview()
        elif extension in ['.pdf']:
            self._preview_pdf(file_path)
            self._show_ai_for_pdf()
        elif extension in ['.doc', '.docx', '.odt']:
            self._preview_word_document(file_path)
            self._show_ai_for_simple_preview()
        elif extension in ['.xls', '.xlsx']:
            self._preview_excel_document(file_path)
        elif extension in ['.ppt', '.pptx']:
            self._preview_powerpoint_document(file_path)
        else:
            self._show_unsupported(extension)
    
    def _reset_preview_state(self):
        """Reset preview state when switching documents"""
        # Clear any text-only UI (action buttons, return button)
        self._clear_text_mode_widgets()
        
        # Hide navigation controls initially
        self._nav_widget.setVisible(False)
        # Disable OCR button until a PDF is loaded
        if hasattr(self, '_ocr_btn'):
            self._ocr_btn.setEnabled(False)
            self._ocr_btn.setVisible(False)
        if hasattr(self, '_ai_btn'):
            self._ai_btn.setEnabled(False)
            self._ai_btn.setVisible(False)

    def _clear_text_mode_widgets(self):
        """Remove text-mode controls like return button and action buttons"""
        if hasattr(self, '_action_buttons_widget'):
            self._action_buttons_widget.setVisible(False)
            self._action_buttons_widget.deleteLater()
            delattr(self, '_action_buttons_widget')
        if hasattr(self, '_return_btn'):
            self._return_btn.setVisible(False)
            self._return_btn.deleteLater()
            delattr(self, '_return_btn')
        if hasattr(self, '_refresh_btn'):
            self._refresh_btn.setVisible(False)
            self._refresh_btn.deleteLater()
            delattr(self, '_refresh_btn')
    
    def _go_to_previous_page(self):
        """Go to the previous page of the PDF"""
        if self.current_page_num > 0:
            self.current_page_num -= 1
            self._render_current_pdf_page()
    
    def _go_to_next_page(self):
        """Go to the next page of the PDF"""
        if self.current_page_num < self.total_pages - 1:
            self.current_page_num += 1
            self._render_current_pdf_page()
    
    def _extract_pdf_text(self):
        """Extract text from the current PDF page and show it in a text window"""
        if not self.current_file_path:
            return
        
        file_info = Path(self.current_file_path)
        extension = file_info.suffix.lower()
        
        if extension == '.pdf' and self.current_pdf_doc:
            # Extract text from PDF
            self._extract_text_from_pdf()
        else:
            self._show_error(self.tr("Text extraction is available only for PDF files"))

    def _analyze_with_ai(self, force_refresh=False):
        """Analyze current file with Gemini and show/propagate results"""
        if not self.current_file_path:
            return

        # Check if we have a cached result and not forcing refresh
        if not force_refresh and self.current_file_path in self._ai_result_cache:
            cached_result = self._ai_result_cache[self.current_file_path]
            self._ensure_textedit_widget()
            self._handle_ai_result(cached_result['result'], cached_result['header'])
            self._add_text_action_buttons()
            self._add_refresh_button()
            self._add_return_button()
            return

        api_key = config.get_gemini_api_key_plain()
        if not api_key:
            self._show_error(self.tr("Set the Gemini API key in Preferences to use AI."))
            return

        prompt = (
            "Analyze the document provided and extract the following information. "
            "You must find and return all fields even if some are missing."
            "The fields to find are: the document date, the sender/organization name, the document subject/topic, and the recipient name."
            "Return ONLY a valid JSON object (no other text) with exactly these fields:\n\n"
            "{\n"
            "  \"ocr_text\": \"The complete OCR/text content of the document, preserving line breaks\",\n"
            "  \"file_date\": \"The document date in DD-MM-YYYY format, or 'None' if not found\",\n"
            "  \"file_organization\": \"The sender/organization name, or 'None' if not found\",\n"
            "  \"file_subject\": \"The document subject/topic, or 'None' if not found\",\n"
            "  \"file_receiver\": \"The recipient name, or 'None' if not found\"\n"
            "}\n\n"
            "Important:\n"
            "- ocr_text MUST contain ONLY the document text content, nothing else\n"
            "- All fields must be present in the JSON\n"
            "- Use 'None' (as a string) for any field that cannot be determined\n"
            "- Return ONLY the JSON, no additional text or explanations"
        )

        # Determine file type and choose appropriate analysis method
        extension = Path(self.current_file_path).suffix.lower()
        office_extensions = ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']
        
        if extension in office_extensions:
            # For Office documents, extract text first then send to Gemini
            extract_success, text_or_error = self._extract_text_from_office(self.current_file_path)
            if not extract_success:
                self._ensure_textedit_widget()
                self._preview_widget.setPlainText(f"❌ {self.tr('Text extraction failed')}: {text_or_error}")
                self._add_text_action_buttons()
                self._add_refresh_button()
                self._add_return_button()
                return
            
            success, result = analyze_text_with_gemini(api_key, text_or_error, prompt)
        else:
            # For PDFs, images, and text files, send the file directly
            success, result = analyze_file_with_gemini(api_key, self.current_file_path, prompt)

        self._ensure_textedit_widget()
        if success:
            header = f"🤖 {self.tr('Analyze with AI')} - {Path(self.current_file_path).name}\n" + "=" * 50 + "\n\n"
            # Cache the result
            self._ai_result_cache[self.current_file_path] = {
                'result': result,
                'header': header
            }
            self._handle_ai_result(result, header)
        else:
            self._preview_widget.setPlainText(f"❌ {self.tr('AI analysis failed')}: {result}")

        self._add_text_action_buttons()
        self._add_refresh_button()
        self._add_return_button()
    
    def _handle_ai_result(self, result_text: str, header: str):
        """Parse AI JSON, show OCR text, and emit field fills."""
        # Clean up the response - remove markdown code blocks if present
        cleaned_text = result_text.strip()
        
        # Remove markdown code block wrapper (```json ... ``` or ``` ... ```)
        if cleaned_text.startswith("```"):
            # Find the first newline after opening ```
            first_newline = cleaned_text.find("\n")
            if first_newline != -1:
                cleaned_text = cleaned_text[first_newline + 1:]
            # Remove closing ```
            if cleaned_text.endswith("```"):
                cleaned_text = cleaned_text[:-3]
            cleaned_text = cleaned_text.strip()
        
        parsed = None
        try:
            parsed = json.loads(cleaned_text)
        except Exception as e:
            print(f"JSON parsing error: {e}")
            print(f"Attempted to parse: {cleaned_text[:200]}...")
            parsed = None

        ocr_text = None
        file_date = None
        file_org = None
        file_subject = None
        file_receiver = None

        if isinstance(parsed, dict):
            ocr_text = parsed.get("ocr_text")
            file_date = parsed.get("file_date")
            file_org = parsed.get("file_organization")
            file_subject = parsed.get("file_subject")
            file_receiver = parsed.get("file_receiver")
        else:
            # If not JSON, just show raw text
            ocr_text = result_text

        # Normalize values
        ocr_text_norm = self._format_ocr_text(ocr_text)
        date_norm = self._normalize_ai_value(file_date)
        org_norm = self._normalize_ai_value(file_org)
        subject_norm = self._normalize_ai_value(file_subject)
        receiver_norm = self._normalize_ai_value(file_receiver)
        #debug print
        print("AI Extracted Values:")
        print(f"Date: {date_norm}")
        print(f"Organization: {org_norm}")
        print(f"Subject: {subject_norm}")
        print(f"Receiver: {receiver_norm}")
        print(f"OCR Text: {ocr_text_norm[:100]}...")  # Print first 100 chars

        # Show only OCR text in preview (no metadata)
        content_to_show = ocr_text_norm if ocr_text_norm else self.tr("No OCR text extracted")
        self._extracted_text_content = header + content_to_show
        self._is_showing_extracted_text = True
        self._preview_widget.setPlainText(self._extracted_text_content)

        # Emit signals to populate renamer fields if present
        if date_norm:
            self.send_to_date_requested.emit(date_norm)
        if org_norm:
            self.send_to_organization_requested.emit(org_norm)
        if subject_norm:
            self.send_to_subject_requested.emit(subject_norm)
        if receiver_norm:
            self.send_to_receiver_requested.emit(receiver_norm)
 
    def _normalize_ai_value(self, value):
        """Return a cleaned string or empty if value is None/"none"/blank."""
        if value is None:
            return ""
        if isinstance(value, (list, tuple)):
            value = " ".join(str(v) for v in value if v)
        value_str = str(value).strip()
        if not value_str:
            return ""
        if value_str.lower() == "none":
            return ""
        return value_str
    
    def _format_ocr_text(self, ocr_text):
        """Format OCR text by converting escaped newlines to actual newlines"""
        if not ocr_text:
            return ""
        if isinstance(ocr_text, (list, tuple)):
            ocr_text = " ".join(str(v) for v in ocr_text if v)
        text_str = str(ocr_text).strip()
        if not text_str:
            return ""
        if text_str.lower() == "none":
            return ""
        # Convert escaped newlines to actual newlines
        text_str = text_str.replace("\\n", "\n")
        return text_str

    def _extract_text_from_pdf(self):
        """Extract text from PDF page"""
        try:
            # Get the current page
            page = self.current_pdf_doc[self.current_page_num]
            
            # Extract text content
            text_content = page.get_text()
            
            if text_content.strip():
                # Show extracted text
                header = f"📄 {self.tr('Extracted Text')} - {self.tr('Page')} {self.current_page_num + 1}\n"
                self._show_extracted_text(text_content, header)
            else:
                self._show_error(self.tr("No text found on this page"))
                
        except Exception as e:
            self._show_error(f"{self.tr('Error extracting text from PDF')}: {str(e)}")
    
    def _show_extracted_text(self, text_content, header_text=None):
        """Show extracted text in a separate window/dialog"""
        # For now, replace the current preview with the text
        # In the future, this could open a separate dialog
        self._ensure_textedit_widget()
        
        # Styling is now applied in _ensure_textedit_widget()
        
        # Add header to indicate this is extracted text
        if header_text is None:
            if self.current_pdf_doc:
                header_text = f"📄 {self.tr('Extracted Text')} - {self.tr('Page')} {self.current_page_num + 1}\n"
            else:
                header_text = f"🔤 {self.tr('Extracted Text')} - {Path(self.current_file_path).name}\n"
        
        header = header_text + "=" * 50 + "\n\n"
        
        self._extracted_text_content = header + text_content
        self._is_showing_extracted_text = True
        self._preview_widget.setPlainText(self._extracted_text_content)
        
        # Add action buttons and return button
        self._add_text_action_buttons()
        self._add_return_button()
    
    def _add_text_action_buttons(self):
        """Add buttons to send selected text to form fields"""
        # Remove existing action buttons if they exist
        if hasattr(self, '_action_buttons_widget'):
            self._action_buttons_widget.setVisible(False)
            self._action_buttons_widget.deleteLater()
            delattr(self, '_action_buttons_widget')
        
        # Create a widget to hold the action buttons
        self._action_buttons_widget = QWidget()
        button_layout = QHBoxLayout(self._action_buttons_widget)
        button_layout.setContentsMargins(5, 5, 5, 5)
        button_layout.setSpacing(10)
        
        # Style for action buttons
        button_style = """
            QPushButton {
                background-color: #007bff;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 8px 12px;
                font-weight: bold;
                font-size: 11px;
            }
            QPushButton:hover {
                background-color: #0056b3;
            }
            QPushButton:pressed {
                background-color: #004085;
            }
            QPushButton:disabled {
                background-color: #6c757d;
                color: #adb5bd;
            }
        """
        
        # Send to Date button
        self._send_to_date_btn = QPushButton("📅 " + self.tr("Send to Date"))
        self._send_to_date_btn.setStyleSheet(button_style)
        self._send_to_date_btn.clicked.connect(self._on_send_to_date)
        button_layout.addWidget(self._send_to_date_btn)
        
        # Send to Organization button
        self._send_to_org_btn = QPushButton("🏢 " + self.tr("Send to Organization"))
        self._send_to_org_btn.setStyleSheet(button_style)
        self._send_to_org_btn.clicked.connect(self._on_send_to_organization)
        button_layout.addWidget(self._send_to_org_btn)
        
        # Send to Subject button
        self._send_to_subject_btn = QPushButton("📝 " + self.tr("Send to Subject"))
        self._send_to_subject_btn.setStyleSheet(button_style)
        self._send_to_subject_btn.clicked.connect(self._on_send_to_subject)
        button_layout.addWidget(self._send_to_subject_btn)
        
        # Send to Receiver button
        self._send_to_receiver_btn = QPushButton("👤 " + self.tr("Send to Receiver"))
        self._send_to_receiver_btn.setStyleSheet(button_style)
        self._send_to_receiver_btn.clicked.connect(self._on_send_to_receiver)
        button_layout.addWidget(self._send_to_receiver_btn)
        
        # Add to main layout at the bottom
        main_layout = self.layout()
        main_layout.addWidget(self._action_buttons_widget)
    
    def _on_send_to_date(self):
        """Handle send to date button click"""
        selected_text = self._get_selected_text()
        if selected_text:
            self.send_to_date_requested.emit(selected_text)
        else:
            self._show_no_selection_message()
    
    def _on_send_to_organization(self):
        """Handle send to organization button click"""
        selected_text = self._get_selected_text()
        if selected_text:
            self.send_to_organization_requested.emit(selected_text)
        else:
            self._show_no_selection_message()
    
    def _on_send_to_subject(self):
        """Handle send to subject button click"""
        selected_text = self._get_selected_text()
        if selected_text:
            self.send_to_subject_requested.emit(selected_text)
        else:
            self._show_no_selection_message()
    
    def _on_send_to_receiver(self):
        """Handle send to receiver button click"""
        selected_text = self._get_selected_text()
        if selected_text:
            self.send_to_receiver_requested.emit(selected_text)
        else:
            self._show_no_selection_message()
    
    def _get_selected_text(self):
        """Get the currently selected text from the preview widget"""
        if isinstance(self._preview_widget, QTextEdit):
            cursor = self._preview_widget.textCursor()
            return cursor.selectedText().strip()
        return ""
    
    def _show_no_selection_message(self):
        """Show message when no text is selected"""
        from PySide6.QtWidgets import QMessageBox
        QMessageBox.information(
            self,
            self.tr("No Text Selected"),
            self.tr("Please select some text first before using this button.")
        )
    
    def _add_return_button(self):
        """Add a button to return to original document view"""
        if hasattr(self, '_return_btn'):
            return  # Button already exists
            
        # Create a return button and add it to the navigation
        if self.current_pdf_doc:
            button_text = "📄 " + self.tr("Back to PDF")
        else:
            button_text = "�️ " + self.tr("Back to Image")
            
        self._return_btn = QPushButton(button_text)
        self._return_btn.setMaximumHeight(30)
        self._return_btn.setStyleSheet("""
            QPushButton {
                background-color: #28a745;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 5px 10px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #218838;
            }
            QPushButton:pressed {
                background-color: #1e7e34;
            }
        """)
        self._return_btn.clicked.connect(self._return_to_original_view)
        
        # Add to navigation layout
        nav_layout = self._nav_widget.layout()
        nav_layout.addWidget(self._return_btn)
    
    def _add_refresh_button(self):
        """Add a refresh button to re-run AI analysis"""
        # Remove existing refresh button if it exists
        if hasattr(self, '_refresh_btn'):
            return  # Button already exists
        
        self._refresh_btn = QPushButton("🔄 " + self.tr("Refresh"))
        self._refresh_btn.setMaximumHeight(30)
        self._refresh_btn.setStyleSheet("""
            QPushButton {
                background-color: #17a2b8;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 5px 10px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #138496;
            }
            QPushButton:pressed {
                background-color: #0c5460;
            }
        """)
        self._refresh_btn.clicked.connect(self._on_refresh_ai_clicked)
        
        # Add to navigation layout
        nav_layout = self._nav_widget.layout()
        nav_layout.addWidget(self._refresh_btn)
    
    def _extract_text_from_office(self, file_path: str) -> tuple[bool, str]:
        """Extract text from Office documents (Word, Excel, PowerPoint)."""
        extension = Path(file_path).suffix.lower()
        
        try:
            # Word documents
            if extension in ['.docx', '.doc']:
                if not DOCX_AVAILABLE:
                    return False, "python-docx not available"
                doc = Document(file_path)
                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                return True, "\n".join(paragraphs)
            
            # Excel documents
            elif extension in ['.xlsx', '.xls']:
                if extension == '.xlsx':
                    if not OPENPYXL_AVAILABLE:
                        return False, "openpyxl not available"
                    wb = load_workbook(file_path, data_only=True)
                    content_parts = []
                    for sheet_name in wb.sheetnames[:5]:  # Limit to first 5 sheets
                        sheet = wb[sheet_name]
                        content_parts.append(f"[Sheet: {sheet_name}]")
                        for row in sheet.iter_rows(max_row=100, values_only=True):
                            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                content_parts.append(row_text)
                    return True, "\n".join(content_parts)
                else:  # .xls
                    if not XLRD_AVAILABLE:
                        return False, "xlrd not available"
                    wb = xlrd.open_workbook(file_path)
                    content_parts = []
                    for sheet_idx in range(min(5, wb.nsheets)):
                        sheet = wb.sheet_by_index(sheet_idx)
                        content_parts.append(f"[Sheet: {sheet.name}]")
                        for row_idx in range(min(100, sheet.nrows)):
                            row_text = " | ".join([str(cell.value) for cell in sheet.row(row_idx)])
                            if row_text.strip():
                                content_parts.append(row_text)
                    return True, "\n".join(content_parts)
            
            # PowerPoint documents
            elif extension in ['.pptx', '.ppt']:
                if not PPTX_AVAILABLE:
                    return False, "python-pptx not available"
                prs = Presentation(file_path)
                content_parts = []
                for i, slide in enumerate(prs.slides[:20], 1):  # Limit to first 20 slides
                    content_parts.append(f"[Slide {i}]")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content_parts.append(shape.text.strip())
                return True, "\n".join(content_parts)
            
            else:
                return False, f"Unsupported file type: {extension}"
                
        except Exception as e:
            return False, f"Failed to extract text: {e}"
    
    def _on_refresh_ai_clicked(self):
        """Handle refresh button click with confirmation"""
        from PySide6.QtWidgets import QMessageBox
        
        reply = QMessageBox.question(
            self,
            self.tr("Refresh AI Analysis"),
            self.tr("This will re-run the AI analysis. Continue?"),
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No
        )
        
        if reply == QMessageBox.Yes:
            self._analyze_with_ai(force_refresh=True)
    
    def _return_to_original_view(self):
        """Return to original document view"""
        # Clean up any text-mode widgets
        self._clear_text_mode_widgets()
        
        # Clear extracted text state
        self._extracted_text_content = None
        self._is_showing_extracted_text = False
        
        # Re-render the current document
        if self.current_file_path:
            file_info = Path(self.current_file_path)
            extension = file_info.suffix.lower()
            
            if extension == '.pdf':
                self._render_current_pdf_page()
            elif extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif']:
                self._preview_image(self.current_file_path)
    
    def _update_navigation_buttons(self):
        """Update the state of navigation buttons"""
        if not PYMUPDF_AVAILABLE or not self.current_pdf_doc:
            self._nav_widget.setVisible(False)
            return
        
        # Always show navigation controls for PDFs (including single page for OCR)
        self._nav_widget.setVisible(True)
        
        # Show all PDF navigation elements
        self._prev_page_btn.setVisible(True)
        self._next_page_btn.setVisible(True)
        self._page_info_label.setVisible(True)
        self._ocr_btn.setVisible(True)
        if hasattr(self, '_ai_btn'):
            self._ai_btn.setVisible(True)
            self._ai_btn.setEnabled(True)
        
        # Update OCR button text for PDFs
        self._ocr_btn.setText("🔤 " + self.tr("Extract Text"))
        
        # Update page navigation buttons
        if self.total_pages > 1:
            self._prev_page_btn.setEnabled(self.current_page_num > 0)
            self._next_page_btn.setEnabled(self.current_page_num < self.total_pages - 1)
            self._page_info_label.setText(f"{self.tr('Page')} {self.current_page_num + 1} {self.tr('of')} {self.total_pages}")
        else:
            self._prev_page_btn.setEnabled(False)
            self._next_page_btn.setEnabled(False)
            self._page_info_label.setText(f"{self.tr('Page')} 1 {self.tr('of')} 1")
        
        # Enable OCR button when PDF is loaded
        self._ocr_btn.setEnabled(True)

    def _show_ai_for_pdf(self):
        """Ensure AI button is shown alongside PDF controls"""
        if hasattr(self, '_ai_btn'):
            self._ai_btn.setVisible(True)
            self._ai_btn.setEnabled(True)
        # Keep PDF navigation state updated
        self._update_navigation_buttons()

    def _show_ai_for_simple_preview(self):
        """Show only the AI button for non-PDF previews (image/text/doc)"""
        if not hasattr(self, '_ai_btn'):
            return
        self._nav_widget.setVisible(True)
        self._prev_page_btn.setVisible(False)
        self._next_page_btn.setVisible(False)
        self._page_info_label.setVisible(False)
        if hasattr(self, '_ocr_btn'):
            self._ocr_btn.setVisible(False)
            self._ocr_btn.setEnabled(False)
        self._ai_btn.setVisible(True)
        self._ai_btn.setEnabled(True)
    
    def _preview_image(self, file_path):
        """Preview image files"""
        try:
            # Ensure we have a QLabel widget for displaying images
            self._ensure_label_widget()
            
            pixmap = QPixmap(file_path)
            if pixmap.isNull():
                self._show_error(self.tr("Cannot load image"))
                return
            
            # Scale image to fit the available space while maintaining aspect ratio
            # Get current widget size, with some padding
            widget_size = self._preview_widget.size()
            available_width = max(widget_size.width() - 40, 300)  # Min 300px
            available_height = max(widget_size.height() - 40, 400)  # Min 400px
            max_size = QSize(available_width, available_height)
            
            scaled_pixmap = pixmap.scaled(max_size, Qt.KeepAspectRatio, Qt.SmoothTransformation)
            
            self._preview_widget.setPixmap(scaled_pixmap)
            self._preview_widget.setText("")
            
            # Show AI button for images (OCR remains disabled)
            self._show_ai_for_simple_preview()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error loading image')}: {str(e)}")
    
    def _preview_text(self, file_path):
        """Preview text files"""
        try:
            # Replace the label with a text edit for better text display
            if not isinstance(self._preview_widget, QTextEdit):
                self._content_layout.removeWidget(self._preview_widget)
                self._preview_widget.deleteLater()
                
                self._preview_widget = QTextEdit()
                self._preview_widget.setReadOnly(True)
                self._preview_widget.setMinimumSize(300, 400)
                self._preview_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
                self._content_layout.addWidget(self._preview_widget)
            
            # Try to detect encoding via BOM first
            detected_encoding = self._detect_text_file_encoding(file_path)

            # Fallback list of encodings to try
            encodings = ([detected_encoding] if detected_encoding else []) + [
                'utf-8', 'utf-8-sig', 'utf-16', 'utf-16-le', 'utf-16-be', 'latin-1', 'cp1252'
            ]
            # Remove duplicates while preserving order
            seen = set()
            encodings = [e for e in encodings if e and not (e in seen or seen.add(e))]

            content = None
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='strict') as file:
                        content = file.read(10000)  # Limit to first 10KB
                        break
                except UnicodeDecodeError:
                    continue
                except Exception:
                    continue
            
            if content is not None:
                if len(content) >= 10000:
                    content += f"\n\n{self.tr('[File truncated - showing first 10KB]')}"
                self._preview_widget.setPlainText(content)
            else:
                self._show_error(self.tr("Cannot decode text file"))
                
        except Exception as e:
            self._show_error(f"{self.tr('Error reading text file')}: {str(e)}")

    def _detect_text_file_encoding(self, file_path):
        """Detect text file encoding using BOM if present"""
        try:
            with open(file_path, 'rb') as f:
                head = f.read(4)
            # BOM checks
            if head.startswith(b'\xef\xbb\xbf'):
                return 'utf-8-sig'
            if head.startswith(b'\xff\xfe\x00\x00'):
                return 'utf-32-le'
            if head.startswith(b'\x00\x00\xfe\xff'):
                return 'utf-32-be'
            if head.startswith(b'\xff\xfe'):
                return 'utf-16-le'
            if head.startswith(b'\xfe\xff'):
                return 'utf-16-be'
        except Exception:
            pass
        return None
    
    def _preview_pdf(self, file_path):
        """Preview PDF files using PyMuPDF"""
        if not PYMUPDF_AVAILABLE:
            self._show_pdf_fallback(file_path)
            return
            
        try:
            # Close previous PDF document if exists
            if self.current_pdf_doc:
                self.current_pdf_doc.close()
                self.current_pdf_doc = None
            
            # Open PDF document
            self.current_pdf_doc = fitz.open(file_path)
            
            if len(self.current_pdf_doc) == 0:
                self._show_error(self.tr("PDF file appears to be empty"))
                self.current_pdf_doc.close()
                self.current_pdf_doc = None
                return
            
            # Set page navigation info
            self.current_page_num = 0  # Start with first page
            self.total_pages = len(self.current_pdf_doc)
            
            # Render the first page
            self._render_current_pdf_page()
            
            # Update navigation buttons
            self._update_navigation_buttons()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error loading PDF')}: {str(e)}")
            if self.current_pdf_doc:
                self.current_pdf_doc.close()
                self.current_pdf_doc = None
    
    def _render_current_pdf_page(self):
        """Render the current page of the PDF as image"""
        if not self.current_pdf_doc or self.current_page_num >= len(self.current_pdf_doc):
            return
            
        # Ensure we're not showing text-mode controls while rendering PDF pages
        self._clear_text_mode_widgets()

        try:
            # Get the current page
            page = self.current_pdf_doc[self.current_page_num]
            
            # Always render as image for consistent preview
            self._render_pdf_as_image(page)
            
            # Update navigation buttons
            self._update_navigation_buttons()
            
            # Add tooltip with page information
            if self.total_pages > 1:
                self._preview_widget.setToolTip(
                    self.tr(f"PDF Document - Page {self.current_page_num + 1} of {self.total_pages}")
                )
            else:
                self._preview_widget.setToolTip(self.tr("PDF Document - 1 page"))
                
        except Exception as e:
            self._show_error(f"{self.tr('Error rendering PDF page')}: {str(e)}")
    
    def _render_pdf_as_image(self, page):
        """Render PDF page as image when no text is available"""
        try:
            # Get page dimensions
            page_rect = page.rect
            
            # Calculate appropriate zoom level based on available widget size
            widget_size = self._preview_widget.size()
            available_width = max(widget_size.width() - 40, 300)
            available_height = max(widget_size.height() - 40, 400)
            
            # Calculate zoom to fit the available space
            zoom_x = available_width / page_rect.width
            zoom_y = available_height / page_rect.height
            zoom = min(zoom_x, zoom_y, 2.0)  # Cap zoom at 2.0 for readability
            
            # Render page to pixmap
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to QPixmap
            img_data = pix.tobytes("ppm")
            qpixmap = QPixmap()
            qpixmap.loadFromData(img_data)
            
            # Clean up pixmap
            pix = None
            
            if qpixmap.isNull():
                self._show_error(self.tr("Failed to render PDF page"))
                return
            
            # Ensure we have a label widget for displaying the pixmap
            self._ensure_label_widget()
            
            # Clear any existing styling and display the PDF
            self._preview_widget.setStyleSheet("")
            self._preview_widget.setPixmap(qpixmap)
                
        except Exception as e:
            self._show_error(f"{self.tr('Error rendering PDF as image')}: {str(e)}")
    
    def _ensure_textedit_widget(self):
        """Ensure the preview widget is a QTextEdit for selectable text"""
        if not isinstance(self._preview_widget, QTextEdit):
            self._content_layout.removeWidget(self._preview_widget)
            self._preview_widget.deleteLater()
            
            self._preview_widget = QTextEdit()
            self._preview_widget.setReadOnly(True)
            self._preview_widget.setMinimumSize(300, 400)
            self._preview_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
            self._preview_widget.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
            self._preview_widget.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
            self._content_layout.addWidget(self._preview_widget)
        
        # Apply consistent styling
        self._preview_widget.setStyleSheet("""
            QTextEdit {
                background-color: white;
                color: black;
                border: 2px solid #007acc;
                border-radius: 8px;
                padding: 15px;
                font-family: 'Segoe UI', Arial, sans-serif;
                font-size: 12pt;
                line-height: 1.5;
                selection-background-color: rgba(0, 120, 215, 0.3);
                selection-color: black;
            }
        """)
    
    def _show_pdf_fallback(self, file_path):
        """Show PDF fallback message when PyMuPDF is not available"""
        self._ensure_label_widget()
        self._preview_widget.setText(
            f"📄 {self.tr('PDF Document')}\n\n"
            f"{self.tr('PDF preview requires PyMuPDF library.')}\n"
            f"{self.tr('Install with: pip install PyMuPDF')}\n\n"
            f"{self.tr('File:')} {os.path.basename(file_path)}"
        )
        self._preview_widget.setStyleSheet("""
            QLabel {
                background-color: #fff3cd;
                border: 2px dashed #ffc107;
                border-radius: 10px;
                padding: 20px;
                color: #856404;
            }
        """)
    
    def _preview_word_document(self, file_path):
        """Preview Word documents (.doc, .docx)"""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.docx' and DOCX_AVAILABLE:
            self._preview_docx_file(file_path)
        elif file_extension == '.doc':
            self._preview_legacy_doc_file(file_path)
        else:
            self._show_office_fallback(file_path, "Word Document", "📄")
    
    def _preview_excel_document(self, file_path):
        """Preview Excel documents (.xls, .xlsx)"""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.xlsx' and OPENPYXL_AVAILABLE:
            self._preview_xlsx_file(file_path)
        elif file_extension == '.xls':
            self._preview_legacy_xls_file(file_path)
        else:
            self._show_office_fallback(file_path, "Excel Spreadsheet", "📊")
    
    def _preview_powerpoint_document(self, file_path):
        """Preview PowerPoint documents (.ppt, .pptx)"""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pptx' and PPTX_AVAILABLE:
            self._preview_pptx_file(file_path)
        elif file_extension == '.ppt':
            self._preview_legacy_ppt_file(file_path)
        else:
            self._show_office_fallback(file_path, "PowerPoint Presentation", "📽️")
    
    def _preview_docx_file(self, file_path):
        """Preview .docx file using python-docx"""
        try:
            doc = Document(file_path)
            
            # Extract text content from paragraphs
            content_parts = []
            paragraph_count = 0
            max_paragraphs = 50  # Limit number of paragraphs to prevent overwhelming UI
            
            for paragraph in doc.paragraphs:
                if paragraph_count >= max_paragraphs:
                    content_parts.append(f"\n... ({self.tr('showing first')} {max_paragraphs} {self.tr('paragraphs')})")
                    break
                    
                text = paragraph.text.strip()
                if text:  # Only add non-empty paragraphs
                    content_parts.append(text)
                    paragraph_count += 1
            
            # Show as text in a QTextEdit
            self._ensure_textedit_widget()
            
            if content_parts:
                header = f"📄 {self.tr('Word Document')}: {Path(file_path).name}\n" + "=" * 60 + "\n\n"
                full_content = header + "\n\n".join(content_parts)
            else:
                full_content = f"📄 {self.tr('Word Document')}: {Path(file_path).name}\n\n{self.tr('No readable text content found.')}"
            
            self._preview_widget.setPlainText(full_content)
            self._preview_widget.setStyleSheet("""
                QTextEdit {
                    background-color: white;
                    color: black;
                    border: 2px solid #28a745;
                    border-radius: 8px;
                    padding: 15px;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 11pt;
                    line-height: 1.4;
                }
            """)
            
            # Add text action buttons for extracted content
            self._add_text_action_buttons()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading Word document')}: {str(e)}")
    
    def _preview_xlsx_file(self, file_path):
        """Preview .xlsx file using openpyxl"""
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            
            content_parts = []
            content_parts.append(f"📊 {self.tr('Excel Spreadsheet')}: {Path(file_path).name}")
            content_parts.append("=" * 60)
            content_parts.append("")
            
            # Limit to first few sheets and rows
            max_sheets = 3
            max_rows = 20
            max_cols = 10
            
            sheet_count = 0
            for sheet_name in workbook.sheetnames:
                if sheet_count >= max_sheets:
                    content_parts.append(f"\n... ({self.tr('showing first')} {max_sheets} {self.tr('sheets')})")
                    break
                
                sheet = workbook[sheet_name]
                content_parts.append(f"\n📋 {self.tr('Sheet')}: {sheet_name}")
                content_parts.append("-" * 40)
                
                # Read sheet data
                row_count = 0
                for row in sheet.iter_rows(max_row=max_rows, max_col=max_cols, values_only=True):
                    if row_count >= max_rows:
                        content_parts.append(f"... ({self.tr('showing first')} {max_rows} {self.tr('rows')})")
                        break
                    
                    # Filter out empty rows and format cells
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            row_data.append(str(cell))
                        else:
                            row_data.append("")
                    
                    # Only add rows that have some content
                    if any(cell.strip() for cell in row_data):
                        content_parts.append(" | ".join(row_data))
                        row_count += 1
                
                sheet_count += 1
            
            workbook.close()
            
            # Show as text in a QTextEdit
            self._ensure_textedit_widget()
            full_content = "\n".join(content_parts)
            
            self._preview_widget.setPlainText(full_content)
            self._preview_widget.setStyleSheet("""
                QTextEdit {
                    background-color: white;
                    color: black;
                    border: 2px solid #17a2b8;
                    border-radius: 8px;
                    padding: 15px;
                    font-family: 'Courier New', monospace;
                    font-size: 10pt;
                    line-height: 1.3;
                }
            """)
            
            # Add text action buttons for extracted content
            self._add_text_action_buttons()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading Excel file')}: {str(e)}")
    
    def _preview_pptx_file(self, file_path):
        """Preview .pptx file using python-pptx"""
        try:
            presentation = Presentation(file_path)
            
            content_parts = []
            content_parts.append(f"�️ {self.tr('PowerPoint Presentation')}: {Path(file_path).name}")
            content_parts.append("=" * 60)
            content_parts.append("")
            
            # Limit number of slides
            max_slides = 10
            slide_count = 0
            
            for slide in presentation.slides:
                if slide_count >= max_slides:
                    content_parts.append(f"\n... ({self.tr('showing first')} {max_slides} {self.tr('slides')})")
                    break
                
                content_parts.append(f"\n🎯 {self.tr('Slide')} {slide_count + 1}")
                content_parts.append("-" * 40)
                
                # Extract text from all shapes in the slide
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    content_parts.extend(slide_texts)
                else:
                    content_parts.append(f"({self.tr('No readable text in this slide')})")
                
                slide_count += 1
            
            # Show as text in a QTextEdit
            self._ensure_textedit_widget()
            full_content = "\n".join(content_parts)
            
            self._preview_widget.setPlainText(full_content)
            self._preview_widget.setStyleSheet("""
                QTextEdit {
                    background-color: white;
                    color: black;
                    border: 2px solid #fd7e14;
                    border-radius: 8px;
                    padding: 15px;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 11pt;
                    line-height: 1.4;
                }
            """)
            
            # Add text action buttons for extracted content
            self._add_text_action_buttons()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading PowerPoint file')}: {str(e)}")
    
    def _preview_legacy_doc_file(self, file_path):
        """Preview legacy .doc file"""
        if WIN32COM_AVAILABLE:
            self._preview_doc_with_com(file_path)
        elif OLEFILE_AVAILABLE:
            self._preview_doc_with_olefile(file_path)
        else:
            self._show_office_fallback(file_path, "Word Document (Legacy)", "📄", 
                                       self.tr("Legacy .doc format requires additional libraries (pywin32 or olefile). Please save as .docx for full preview."))
    
    def _preview_legacy_xls_file(self, file_path):
        """Preview legacy .xls file"""
        if XLRD_AVAILABLE:
            self._preview_xls_with_xlrd(file_path)
        elif WIN32COM_AVAILABLE:
            self._preview_xls_with_com(file_path)
        else:
            self._show_office_fallback(file_path, "Excel Spreadsheet (Legacy)", "📊",
                                       self.tr("Legacy .xls format requires additional libraries (xlrd or pywin32). Please save as .xlsx for full preview."))
    
    def _preview_legacy_ppt_file(self, file_path):
        """Preview legacy .ppt file"""
        if WIN32COM_AVAILABLE:
            self._preview_ppt_with_com(file_path)
        else:
            self._show_office_fallback(file_path, "PowerPoint Presentation (Legacy)", "📽️",
                                       self.tr("Legacy .ppt format requires additional libraries (pywin32). Please save as .pptx for full preview."))
    
    def _preview_doc_with_com(self, file_path):
        """Preview .doc file using Windows COM"""
        try:
            import pythoncom
            pythoncom.CoInitialize()
            
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            
            doc = word.Documents.Open(str(Path(file_path).absolute()))
            
            # Extract text content
            content = doc.Content.Text
            
            # Close document and Word
            doc.Close(False)
            word.Quit()
            pythoncom.CoUninitialize()
            
            # Show as text
            self._show_extracted_office_text(content, file_path, "Word Document (Legacy)", "📄")
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading legacy Word document')}: {str(e)}")
    
    def _preview_doc_with_olefile(self, file_path):
        """Preview .doc file using olefile (basic text extraction)"""
        try:
            if not olefile.isOleFile(file_path):
                self._show_error(self.tr("File is not a valid OLE document"))
                return
                
            # This is a very basic implementation - olefile doesn't directly extract text
            # It can identify the structure but extracting formatted text is complex
            ole = olefile.OleFileIO(file_path)
            
            # Get basic file info
            file_info = Path(file_path)
            content_parts = [
                f"📄 {self.tr('Word Document (Legacy)')}: {file_info.name}",
                "=" * 60,
                "",
                f"{self.tr('File detected as OLE compound document')}",
                f"{self.tr('Size:')} {self._format_file_size(file_info.stat().st_size)}",
                "",
                f"{self.tr('Available streams:')}",
            ]
            
            # List available streams
            for stream in ole.listdir():
                content_parts.append(f"  - {'/'.join(stream)}")
            
            content_parts.extend([
                "",
                f"{self.tr('Note: Full text extraction from .doc files requires Microsoft Word or more advanced libraries.')}"
            ])
            
            ole.close()
            
            self._ensure_textedit_widget()
            self._preview_widget.setPlainText("\n".join(content_parts))
            self._preview_widget.setStyleSheet("""
                QTextEdit {
                    background-color: #fff3cd;
                    color: #856404;
                    border: 2px solid #ffc107;
                    border-radius: 8px;
                    padding: 15px;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 11pt;
                }
            """)
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading legacy Word document')}: {str(e)}")
    
    def _preview_xls_with_xlrd(self, file_path):
        """Preview .xls file using xlrd"""
        try:
            workbook = xlrd.open_workbook(file_path)
            
            content_parts = []
            content_parts.append(f"📊 {self.tr('Excel Spreadsheet (Legacy)')}: {Path(file_path).name}")
            content_parts.append("=" * 60)
            content_parts.append("")
            
            # Limit sheets and rows
            max_sheets = 3
            max_rows = 20
            max_cols = 10
            
            for sheet_idx, sheet_name in enumerate(workbook.sheet_names()[:max_sheets]):
                sheet = workbook.sheet_by_name(sheet_name)
                content_parts.append(f"\n📋 {self.tr('Sheet')}: {sheet_name}")
                content_parts.append("-" * 40)
                
                rows_shown = 0
                for row_idx in range(min(sheet.nrows, max_rows)):
                    row_data = []
                    for col_idx in range(min(sheet.ncols, max_cols)):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.ctype == xlrd.XL_CELL_EMPTY:
                            row_data.append("")
                        elif cell.ctype == xlrd.XL_CELL_TEXT:
                            row_data.append(str(cell.value))
                        elif cell.ctype == xlrd.XL_CELL_NUMBER:
                            row_data.append(str(cell.value))
                        elif cell.ctype == xlrd.XL_CELL_DATE:
                            row_data.append(str(cell.value))
                        else:
                            row_data.append(str(cell.value))
                    
                    # Only add rows with content
                    if any(cell.strip() for cell in row_data if isinstance(cell, str)):
                        content_parts.append(" | ".join(row_data))
                        rows_shown += 1
                
                if sheet.nrows > max_rows:
                    content_parts.append(f"... ({self.tr('showing first')} {max_rows} {self.tr('rows')})")
            
            if len(workbook.sheet_names()) > max_sheets:
                content_parts.append(f"\n... ({self.tr('showing first')} {max_sheets} {self.tr('sheets')})")
            
            self._ensure_textedit_widget()
            self._preview_widget.setPlainText("\n".join(content_parts))
            self._preview_widget.setStyleSheet("""
                QTextEdit {
                    background-color: white;
                    color: black;
                    border: 2px solid #17a2b8;
                    border-radius: 8px;
                    padding: 15px;
                    font-family: 'Courier New', monospace;
                    font-size: 10pt;
                }
            """)
            
            # Add text action buttons
            self._add_text_action_buttons()
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading legacy Excel file')}: {str(e)}")
    
    def _preview_xls_with_com(self, file_path):
        """Preview .xls file using Windows COM"""
        try:
            import pythoncom
            pythoncom.CoInitialize()
            
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            
            workbook = excel.Workbooks.Open(str(Path(file_path).absolute()))
            
            content_parts = []
            content_parts.append(f"📊 {self.tr('Excel Spreadsheet (Legacy)')}: {Path(file_path).name}")
            content_parts.append("=" * 60)
            
            max_sheets = 3
            max_rows = 20
            
            for sheet_idx, sheet in enumerate(workbook.Worksheets[:max_sheets]):
                content_parts.append(f"\n� {self.tr('Sheet')}: {sheet.Name}")
                content_parts.append("-" * 40)
                
                # Get used range
                used_range = sheet.UsedRange
                if used_range is not None:
                    rows = min(used_range.Rows.Count, max_rows)
                    for row_idx in range(1, rows + 1):
                        row_data = []
                        for col_idx in range(1, min(used_range.Columns.Count + 1, 11)):  # Max 10 columns
                            cell_value = sheet.Cells(row_idx, col_idx).Value
                            row_data.append(str(cell_value) if cell_value is not None else "")
                        
                        if any(cell.strip() for cell in row_data):
                            content_parts.append(" | ".join(row_data))
                    
                    if used_range.Rows.Count > max_rows:
                        content_parts.append(f"... ({self.tr('showing first')} {max_rows} {self.tr('rows')})")
            
            workbook.Close(False)
            excel.Quit()
            pythoncom.CoUninitialize()
            
            self._show_extracted_office_text("\n".join(content_parts), file_path, "Excel Spreadsheet (Legacy)", "📊")
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading legacy Excel file')}: {str(e)}")
    
    def _preview_ppt_with_com(self, file_path):
        """Preview .ppt file using Windows COM"""
        try:
            import pythoncom
            pythoncom.CoInitialize()
            
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(str(Path(file_path).absolute()))
            
            content_parts = []
            content_parts.append(f"📽️ {self.tr('PowerPoint Presentation (Legacy)')}: {Path(file_path).name}")
            content_parts.append("=" * 60)
            
            max_slides = 10
            slide_count = min(presentation.Slides.Count, max_slides)
            
            for slide_idx in range(1, slide_count + 1):
                slide = presentation.Slides(slide_idx)
                content_parts.append(f"\n🎯 {self.tr('Slide')} {slide_idx}")
                content_parts.append("-" * 40)
                
                # Extract text from shapes
                slide_texts = []
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        text = shape.TextFrame.TextRange.Text.strip()
                        if text:
                            slide_texts.append(text)
                
                if slide_texts:
                    content_parts.extend(slide_texts)
                else:
                    content_parts.append(f"({self.tr('No readable text in this slide')})")
            
            if presentation.Slides.Count > max_slides:
                content_parts.append(f"\n... ({self.tr('showing first')} {max_slides} {self.tr('slides')})")
            
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            
            self._show_extracted_office_text("\n".join(content_parts), file_path, "PowerPoint Presentation (Legacy)", "📽️")
            
        except Exception as e:
            self._show_error(f"{self.tr('Error reading legacy PowerPoint file')}: {str(e)}")
    
    def _show_extracted_office_text(self, content, file_path, doc_type, icon):
        """Show extracted text from Office documents"""
        self._ensure_textedit_widget()
        
        header = f"{icon} {self.tr(doc_type)}: {Path(file_path).name}\n" + "=" * 60 + "\n\n"
        full_content = header + content
        
        self._preview_widget.setPlainText(full_content)
        self._preview_widget.setStyleSheet("""
            QTextEdit {
                background-color: white;
                color: black;
                border: 2px solid #28a745;
                border-radius: 8px;
                padding: 15px;
                font-family: 'Segoe UI', Arial, sans-serif;
                font-size: 11pt;
                line-height: 1.4;
            }
        """)
        
        # Add text action buttons
        self._add_text_action_buttons()
    
    def _show_office_fallback(self, file_path, doc_type, icon, additional_message=None):
        """Show fallback message for Office documents when libraries aren't available"""
        self._ensure_label_widget()
        
        file_info = Path(file_path)
        message_parts = [
            f"{icon} {self.tr(doc_type)}",
            "",
            f"{self.tr('File:')} {file_info.name}",
            f"{self.tr('Size:')} {self._format_file_size(file_info.stat().st_size)}",
            ""
        ]
        
        if additional_message:
            message_parts.append(additional_message)
        else:
            message_parts.extend([
                f"{self.tr('Preview requires additional libraries.')}",
                f"{self.tr('Install Office document support to view contents.')}"
            ])
        
        self._preview_widget.setText("\n".join(message_parts))
        self._preview_widget.setStyleSheet("""
            QLabel {
                background-color: #fff3cd;
                border: 2px dashed #ffc107;
                border-radius: 10px;
                padding: 20px;
                color: #856404;
            }
        """)
    
    def _format_file_size(self, size_bytes):
        """Format file size in human readable format"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.1f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"
    
    def _show_unsupported(self, extension):
        """Show unsupported file type message"""
        self._ensure_label_widget()
        self._preview_widget.setText(
            f"❓ {self.tr('Unsupported Format')}\n\n"
            f"{self.tr('Cannot preview')} {extension.upper()} {self.tr('files')}"
        )
        self._preview_widget.setStyleSheet("""
            QLabel {
                background-color: #fff3cd;
                border: 2px dashed #ffc107;
                border-radius: 10px;
                padding: 20px;
                color: #856404;
            }
        """)
    
    def _show_error(self, message):
        """Show error message"""
        self._ensure_label_widget()
        self._preview_widget.setText(f"❌ {self.tr('Error')}\n\n{message}")
        self._preview_widget.setStyleSheet("""
            QLabel {
                background-color: #f8d7da;
                border: 2px dashed #dc3545;
                border-radius: 10px;
                padding: 20px;
                color: #721c24;
            }
        """)
    
    def _show_placeholder(self):
        """Show placeholder when no file is selected"""
        self._ensure_label_widget()
        self._preview_widget.setText(
            f"📄 {self.tr('Document Preview')}\n\n"
            f"{self.tr('Select a document from the list to preview it here')}"
        )
        self._preview_widget.setStyleSheet("""
            QLabel {
                background-color: #e9ecef;
                border: 2px dashed #6c757d;
                border-radius: 10px;
                padding: 20px;
                color: #495057;
            }
        """)
    
    def _ensure_label_widget(self):
        """Ensure the preview widget is a QLabel (not QTextEdit)"""
        if not isinstance(self._preview_widget, QLabel):
            self._content_layout.removeWidget(self._preview_widget)
            self._preview_widget.deleteLater()
            
            self._preview_widget = QLabel()
            self._preview_widget.setAlignment(Qt.AlignCenter)
            self._preview_widget.setMinimumSize(300, 400)
            self._preview_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
            self._preview_widget.setWordWrap(True)
            self._preview_widget.setScaledContents(False)
            self._content_layout.addWidget(self._preview_widget)
    
    def _format_file_size(self, size_bytes):
        """Format file size in human readable format"""
        if size_bytes == 0:
            return "0 B"
        
        size_names = ["B", "KB", "MB", "GB"]
        i = 0
        size = float(size_bytes)
        
        while size >= 1024.0 and i < len(size_names) - 1:
            size /= 1024.0
            i += 1
        
        return f"{size:.1f} {size_names[i]}"
    
    def clear_preview(self):
        """Clear the current preview"""
        # Close PDF document if open
        if self.current_pdf_doc:
            self.current_pdf_doc.close()
            self.current_pdf_doc = None
        
        # Reset navigation state
        self.current_page_num = 0
        self.total_pages = 0
        self._nav_widget.setVisible(False)
        
        # Clear extracted text state
        self._extracted_text_content = None
        self._is_showing_extracted_text = False
        self._current_ai_header = None
        
        # Clear other state
        self.current_file_path = None
        self._file_name_label.setText(self.tr("No document selected"))
        self._file_info_label.setText("")
        self._show_placeholder()
    
    def retranslate_ui(self):
        """Retranslate all UI elements"""
        # Update button texts
        if hasattr(self, '_prev_page_btn'):
            self._prev_page_btn.setText("◀ " + self.tr("Previous"))
        if hasattr(self, '_next_page_btn'):
            self._next_page_btn.setText(self.tr("Next") + " ▶")
        if hasattr(self, '_ocr_btn'):
            # OCR only applies to PDFs
            self._ocr_btn.setText("🔤 " + self.tr("Extract Text"))
        if hasattr(self, '_ai_btn'):
            self._ai_btn.setText("🤖 " + self.tr("Analyze with AI"))
        if hasattr(self, '_return_btn'):
            if self.current_pdf_doc:
                self._return_btn.setText("📄 " + self.tr("Back to PDF"))
            else:
                self._return_btn.setText("🖼️ " + self.tr("Back to Image"))
        
        # Update navigation if PDF is loaded
        if self.current_pdf_doc and self.total_pages > 1:
            self._page_info_label.setText(f"{self.tr('Page')} {self.current_page_num + 1} {self.tr('of')} {self.total_pages}")
        elif self.current_pdf_doc:
            self._page_info_label.setText(f"{self.tr('Page')} 1 {self.tr('of')} 1")
        
        # Update preview content
        if self.current_file_path:
            # Re-preview the current file to update language
            self.preview_file(self.current_file_path)
        else:
            self._file_name_label.setText(self.tr("No document selected"))
            self._show_placeholder()
    
    def resizeEvent(self, event):
        """Handle resize events to rescale images and PDFs"""
        super().resizeEvent(event)
        
        # Don't re-render if we're showing extracted text (OCR or AI result)
        if self._is_showing_extracted_text:
            return
        
        # If we're currently previewing content, re-render it for the new size
        if self.current_file_path:
            file_info = Path(self.current_file_path)
            extension = file_info.suffix.lower()
            
            # Check if it's an image and we have a QLabel widget
            if (extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif'] and
                isinstance(self._preview_widget, QLabel) and
                self._preview_widget.pixmap() and not self._preview_widget.pixmap().isNull()):
                # Rescale the image for the new size
                self._preview_image(self.current_file_path)
            
            # Check if it's a PDF
            elif (extension == '.pdf' and PYMUPDF_AVAILABLE and self.current_pdf_doc):
                # Re-render the current PDF page for the new size
                self._render_current_pdf_page()

    def tr(self, text):
        """Translation method - uses parent's tr if available"""
        if self.parent() and hasattr(self.parent(), 'tr'):
            return self.parent().tr(text)
        return text